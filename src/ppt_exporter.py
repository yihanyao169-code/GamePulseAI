from __future__ import annotations

import re
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
from statistics import mean

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.models import AnalysisResult
from src.reporting import sort_category_counts
from src.evaluation import build_single_market_report
from src.methodology import category_standards, sentiment_standards
from src import theme


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


PRIMARY = _rgb("#F5E400")
SECONDARY = _rgb("#333333")
ACCENT = _rgb("#111111")
BACKGROUND = _rgb("#F5F5F2")
CARD = _rgb("#FFFFFF")
TEXT = _rgb("#111111")
MUTED = _rgb("#555555")
BORDER = _rgb("#D9D9D4")
NAVY = _rgb("#F0F0EC")
GREEN = _rgb("#666666")
AMBER = _rgb("#B89F00")
GRID = _rgb("#E6E6E0")
CHART_PALETTE = ["#F5C400", "#2F6BFF", "#27C4C2", "#42B866", "#8A5CF6", "#E84A8A", "#F28C38", "#A5A7AC"]
FONT = "Microsoft YaHei"
REPORT_TIME = datetime.now().strftime("%Y-%m-%d %H:%M")


def get_contrast_text_color(background_hex: str) -> RGBColor:
    value = background_hex.lstrip("#")
    r, g, b = (int(value[index:index + 2], 16) for index in (0, 2, 4))
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return TEXT if luminance > 0.58 else _rgb("#FFFFFF")


def _contrast_text_rgb(fill: RGBColor) -> RGBColor:
    luminance = (0.299 * fill[0] + 0.587 * fill[1] + 0.114 * fill[2]) / 255
    return TEXT if luminance > 0.58 else _rgb("#FFFFFF")


def export_ppt(
    result: AnalysisResult,
    package_name: str,
    bar_chart_path: Path,
    pie_chart_path: Path,
    output_path: Path,
    report: dict | None = None,
    context: dict | None = None,
) -> Path:
    context = _normalize_context(package_name, context)
    rows = [_single_market_row(result, context)]
    report = report or build_single_market_report(result, len(result.classified_reviews))
    presentation = _new_presentation()
    _add_cover(presentation, context)
    _add_score_summary_slide(presentation, rows, report, context)
    _add_category_stats_slide(presentation, result, pie_chart_path, bar_chart_path)
    _add_category_insights_slide(presentation, report)
    _add_strengths_pains_slide(presentation, report)
    _add_sentiment_reviews_slide(presentation, result)
    _add_structured_recommendations_slide(presentation, report)
    _add_methodology_slide(presentation, context)
    _finalize(presentation)
    presentation.save(output_path)
    return output_path


def export_market_comparison_ppt(
    package_name: str,
    comparison_rows: list[dict],
    summary: str,
    chart_paths: list[Path],
    output_path: Path,
    market_results: list[dict] | None = None,
) -> Path:
    presentation = _new_presentation()
    results = [item["result"] for item in market_results or [] if item.get("result")]
    market_label = f"{len(comparison_rows)} Markets" if len(comparison_rows) != 1 else "1 Market"
    _add_cover(presentation, _normalize_context(package_name, {"market": f"跨市场/区域 · {market_label}"}))
    _add_executive_summary(presentation, package_name, comparison_rows, summary, results, "跨市场/区域分析")
    _add_overall_dashboard(presentation, comparison_rows)
    _add_market_comparison_pages(presentation, comparison_rows)
    _add_key_insights(presentation, comparison_rows, summary, results)
    _add_representative_reviews(presentation, comparison_rows, results)
    _add_recommendations(presentation, summary)
    _finalize(presentation)
    presentation.save(output_path)
    return output_path


def _new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    return presentation


def _blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND
    _hud_grid(slide)
    return slide


def _add_cover(presentation: Presentation, context: dict) -> None:
    slide = _blank_slide(presentation)
    game_title = context["game_title"]
    _rect(slide, 0.48, 0.42, 12.38, 6.45, CARD, radius=False, transparency=8, line=BORDER)
    _corner_brackets(slide, 0.48, 0.42, 12.38, 6.45, PRIMARY)
    _text(slide, game_title, 0.82, 0.96, 7.45, 0.72, 36, TEXT, bold=True)
    _text(slide, "Google Play Market Intelligence Report", 0.86, 1.82, 6.8, 0.32, 18, TEXT, bold=True)
    _text(slide, "AI 驱动的移动游戏用户反馈分析与跨市场洞察", 0.86, 2.24, 6.7, 0.28, 15, MUTED)
    _cover_info(slide, "游戏名称", game_title, 0.84, 3.78)
    _cover_info(slide, "Google Play 包名", context["package_name"], 0.84, 4.28)
    _cover_info(slide, "分析市场", context["market"], 0.84, 4.78)
    _cover_info(slide, "评论语言", context["language"], 0.84, 5.28)
    _cover_info(slide, "评论时间", context["time_range"], 0.84, 5.78)
    _cover_info(slide, "分析时间", context["analysis_time"], 0.84, 6.28)
    _text(slide, "AI GENERATED / REPORT READY", 8.72, 5.84, 3.35, 0.22, 12, TEXT, bold=True, align=PP_ALIGN.CENTER)
    _mini_dashboard(slide, 8.6, 1.32)


def _add_executive_summary(
    presentation: Presentation,
    package_name: str,
    rows: list[dict],
    summary: str,
    results: list[AnalysisResult],
    scope: str,
) -> None:
    slide = _content_slide(presentation, "Executive Summary", "核心结论")
    total_comments = _total_comments(rows, results)
    avg_score = _average_score(rows)
    overview = [
        ("分析市场", str(max(len(rows), 1))),
        ("评论数量", str(total_comments)),
        ("平均评分", avg_score),
        ("分析耗时", "--"),
    ]
    _card(slide, 0.65, 1.22, 3.05, 4.65)
    _text(slide, "Overview", 0.95, 1.52, 2.2, 0.32, 19, TEXT, bold=True)
    for index, (label, value) in enumerate(overview):
        y = 2.05 + index * 0.78
        _text(slide, label, 0.95, y, 1.4, 0.22, 11, MUTED, bold=True)
        _text(slide, value, 2.25, y - 0.08, 1.0, 0.34, 22, AMBER, bold=True, align=PP_ALIGN.RIGHT)

    _card(slide, 4.05, 1.22, 8.62, 3.85)
    _text(slide, "Key Findings", 4.38, 1.52, 3.5, 0.32, 19, TEXT, bold=True)
    findings = _findings(summary, results, rows, 5)
    for index, finding in enumerate(findings, start=1):
        y = 2.02 + (index - 1) * 0.53
        _number_dot(slide, index, 4.42, y + 0.03)
        _text(slide, finding, 4.82, y, 7.35, 0.34, 14, TEXT)

    _rect(slide, 0.65, 6.08, 12.02, 0.78, NAVY, line=BORDER)
    _text(slide, "Management Summary", 0.95, 6.27, 2.2, 0.23, 12, TEXT, bold=True)
    _text(slide, _management_summary(summary, rows), 3.0, 6.18, 9.2, 0.38, 15, TEXT, bold=True)


def _add_overall_dashboard(presentation: Presentation, rows: list[dict]) -> None:
    slide = _content_slide(presentation, "Overall Dashboard", "核心经营看板")
    metrics = _dashboard_metrics(rows)
    cards = [
        ("负面率", metrics["negative"], ACCENT),
        ("Top Category", metrics["top_category"], PRIMARY),
        ("BUG占比", metrics["bug"], SECONDARY),
        ("氪金占比", metrics["monetization"], AMBER),
    ]
    positions = [(0.65, 1.28), (6.85, 1.28), (0.65, 4.05), (6.85, 4.05)]
    for (title, value, color), (x, y) in zip(cards, positions):
        _dashboard_tile(slide, title, value, x, y, 5.8, 2.15, color)
    _text(slide, "Dashboard 指标基于本次抓取与分类样本，适合判断短期市场反馈方向。", 0.75, 6.72, 11.5, 0.18, 10, MUTED)


def _add_score_summary_slide(presentation: Presentation, rows: list[dict], report: dict, context: dict | None = None) -> None:
    subtitle = "综合评分与管理层摘要"
    if context:
        subtitle = f"{context['game_title']} · {context['market']} · {context['language']} · {context['time_range']}"
    slide = _content_slide(presentation, "AI Summary", subtitle)
    _card(slide, 0.65, 1.25, 3.2, 2.25)
    if report.get("evaluation_available", True):
        score_text = f"{float(report.get('overall_score', 0)):.1f} / 100"
        grade_text = f"{report.get('grade', '暂无')} · {report.get('grade_label', '')}"
        confidence_text = f"可信度：{report.get('confidence_factor', 0):.2f}"
    else:
        score_text = "需重新分析"
        grade_text = "旧版标注"
        confidence_text = "缺少 v2.0 字段"
    _text(slide, score_text, 0.95, 1.65, 2.4, 0.48, 30, TEXT, bold=True)
    _rect(slide, 0.95, 2.38, 0.9, 0.34, PRIMARY, line=PRIMARY)
    _text(slide, grade_text, 1.03, 2.47, 0.74, 0.1, 9, TEXT, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, confidence_text, 2.05, 2.44, 1.3, 0.18, 12, TEXT, bold=True)
    _text(slide, str(report.get("score_reason", "")), 0.95, 2.92, 2.45, 0.3, 11, MUTED)
    _card(slide, 4.1, 1.25, 8.55, 2.25)
    _text(slide, "整体概括", 4.38, 1.52, 1.4, 0.25, 17, TEXT, bold=True)
    _text(slide, str(report.get("overall_summary", "暂无总结")), 4.38, 1.98, 7.8, 0.85, 14, TEXT)
    metrics = _dashboard_metrics(rows)
    for index, (label, value) in enumerate([("负面率", metrics["negative"]), ("Top Category", metrics["top_category"]), ("BUG占比", metrics["bug"]), ("氪金占比", metrics["monetization"])]):
        x = 0.65 + index * 3.05
        _card(slide, x, 4.1, 2.75, 1.55)
        _text(slide, label, x + 0.2, 4.35, 1.2, 0.18, 11, MUTED, bold=True)
        _text(slide, value, x + 0.2, 4.78, 2.0, 0.35, 22, TEXT, bold=True)
    analyzed_count = int(report.get("analyzed_count", 0) or 0)
    failed_count = int(report.get("failed_count", 0) or 0)
    sample_note = f"本次分析成功处理 {analyzed_count} 条评论。" if not failed_count else f"本次分析成功处理 {analyzed_count} / {analyzed_count + failed_count} 条评论。"
    _text(slide, sample_note + " v2.0：AI 只负责标注，Overall Score 由固定公式计算；不等同于 Google Play 官方评分。", 0.75, 6.62, 11.7, 0.18, 9, MUTED)


def _add_category_stats_slide(presentation: Presentation, result: AnalysisResult, pie_chart_path: Path, bar_chart_path: Path) -> None:
    slide = _content_slide(presentation, "Category Statistics", "类别统计、占比与评论数量排名")
    sorted_counts = sort_category_counts(result.category_counts)
    total = max(sum(sorted_counts.values()), 1)
    _card(slide, 0.65, 1.22, 4.0, 5.55)
    _text(slide, "Rank  Category  Count  Percent", 0.9, 1.55, 3.3, 0.18, 10, MUTED, bold=True)
    for index, (category, count) in enumerate(list(sorted_counts.items())[:11], start=1):
        y = 1.9 + (index - 1) * 0.38
        _text(slide, f"{index:02d}", 0.9, y, 0.35, 0.16, 9, MUTED, bold=True)
        _text(slide, category, 1.32, y, 1.1, 0.16, 10, TEXT, bold=True)
        _text(slide, str(count), 2.65, y, 0.45, 0.16, 10, TEXT, align=PP_ALIGN.RIGHT)
        _text(slide, f"{count / total * 100:.1f}%", 3.4, y, 0.55, 0.16, 10, TEXT, align=PP_ALIGN.RIGHT)
    slide.shapes.add_picture(str(pie_chart_path), Inches(4.95), Inches(1.35), width=Inches(3.65), height=Inches(2.75))
    slide.shapes.add_picture(str(bar_chart_path), Inches(8.75), Inches(1.35), width=Inches(3.65), height=Inches(4.65))


def _add_category_insights_slide(presentation: Presentation, report: dict) -> None:
    slide = _content_slide(presentation, "Category Analysis", "产品团队可执行的类别解读")
    _bullet_cards(slide, report.get("category_insights", []), 0.75, 1.35, 11.8, 5.25, columns=2)


def _add_strengths_pains_slide(presentation: Presentation, report: dict) -> None:
    slide = _content_slide(presentation, "Strengths & Pain Points", "玩家认可的优点与主要痛点")
    _text(slide, "玩家认可的优点", 0.8, 1.25, 3.0, 0.25, 18, TEXT, bold=True)
    _text(slide, "主要痛点", 6.9, 1.25, 3.0, 0.25, 18, TEXT, bold=True)
    _bullet_cards(slide, report.get("strengths", []), 0.75, 1.72, 5.65, 4.95, columns=1)
    _bullet_cards(slide, report.get("pain_points", []), 6.85, 1.72, 5.65, 4.95, columns=1)


def _add_sentiment_reviews_slide(presentation: Presentation, result: AnalysisResult) -> None:
    report = build_single_market_report(result, len(result.classified_reviews))
    slide = _content_slide(presentation, "Sentiment & Reviews", "情感分析与代表性评论")
    sentiments = _sentiment_mix(result)
    for index, key in enumerate(["正面", "中性", "负面"]):
        x = 0.75 + index * 2.0
        _card(slide, x, 1.25, 1.7, 1.1)
        _text(slide, key, x + 0.18, 1.46, 0.9, 0.16, 11, MUTED, bold=True)
        _text(slide, _pct_text(sentiments[key]), x + 0.18, 1.78, 1.1, 0.28, 19, TEXT, bold=True)
    _text(slide, report["sentiment_conclusion"], 0.78, 2.62, 5.55, 0.5, 13, TEXT)
    reviews = result.classified_reviews[:4]
    for index, review in enumerate(reviews):
        x = 6.75
        y = 1.25 + index * 1.28
        _card(slide, x, y, 5.65, 1.0)
        _text(slide, f"{review.category} · {review.sentiment}", x + 0.18, y + 0.16, 1.7, 0.15, 9, TEXT, bold=True)
        _text(slide, _truncate(review.content, 70), x + 0.18, y + 0.4, 5.0, 0.2, 9, TEXT)
        _text(slide, _truncate(review.reason or "暂无概括", 54), x + 0.18, y + 0.68, 5.0, 0.16, 8, MUTED)


def _add_structured_recommendations_slide(presentation: Presentation, report: dict) -> None:
    slide = _content_slide(presentation, "Actionable Recommendations", "P0 / P1 / P2 优化建议")
    recommendations = report.get("recommendations", {})
    for index, level in enumerate(["P0", "P1", "P2"]):
        x = 0.75 + index * 4.05
        _text(slide, {"P0": "P0 立即处理", "P1": "P1 近期优化", "P2": "P2 长期建设"}[level], x, 1.25, 2.2, 0.25, 17, TEXT, bold=True)
        for item_index, item in enumerate(recommendations.get(level, [])[:2]):
            y = 1.78 + item_index * 2.2
            _card(slide, x, y, 3.55, 1.85)
            _text(slide, item.get("title", "建议"), x + 0.18, y + 0.18, 3.0, 0.2, 12, TEXT, bold=True)
            _text(slide, "依据：" + item.get("basis", ""), x + 0.18, y + 0.52, 3.05, 0.25, 9, MUTED)
            _text(slide, "动作：" + item.get("action", ""), x + 0.18, y + 0.92, 3.05, 0.33, 9, MUTED)
            _text(slide, "收益：" + item.get("impact", ""), x + 0.18, y + 1.4, 3.05, 0.22, 9, MUTED)


def _add_methodology_slide(presentation: Presentation, context: dict | None = None) -> None:
    slide = _content_slide(presentation, "Methodology", "方法与口径说明")
    sections = []
    if context:
        sections.append({"title": "本次报告范围", "detail": f"游戏：{context['game_title']}；市场：{context['market']}；评论语言：{context['language']}；评论时间：{context['time_range']}。"})
    sections.extend([
        {"title": "数据来源", "detail": "通过 google-play-scraper 按应用包名、Google Play 商店地区和评论语言请求评论；商店地区不代表评论者国籍或 IP。"},
        {"title": "数据清洗", "detail": "规范空白字符，删除长度小于 2 的评论，正文截断到 1200 字，并按小写正文去重。"},
        {"title": "类别分类", "detail": "Claude 为每条评论选择一个最主要类别；整体评价用于无具体维度的整体态度，其他仅用于无法判断内容。"},
        {"title": "情感判断", "detail": "情感由 Claude 根据评论正文语义判断为正面、中性或负面；当前星级不参与情感分类。"},
        {"title": "AI 严重度标注", "detail": "风险评论标注 S1/S2/S3/S4 和 is_blocking；正面评论标注优势信号。"},
        {"title": "确定性评分", "detail": "Base Score = 55%满意度 + 45%健康度 + 最高+5浮动优势加成；Overall Score = Base Score × 可信度系数。"},
        {"title": "可信度与S4", "detail": "可信度由样本量、分类成功率和泛化评论占比计算；S4采用连续扣分并设置阻塞最低扣分。"},
        {"title": "数据限制", "detail": "报告仅反映本次抓取样本；严重度尚未人工标注集校准，健康度系数为初始经验参数。"},
    ])
    _bullet_cards(slide, sections, 0.75, 1.25, 11.8, 5.45, columns=2)


def _add_market_comparison_pages(presentation: Presentation, rows: list[dict]) -> None:
    chunks = [rows[:4]]
    if len(rows) > 4:
        chunks.append(rows[4:8])
    for index, chunk in enumerate(chunks, start=1):
        suffix = f" ({index}/{len(chunks)})" if len(chunks) > 1 else ""
        slide = _content_slide(presentation, "Market Comparison" + suffix, "市场关键指标横向对比")
        _comparison_quadrant(slide, chunk, "负面率", "负面占比", 0.65, 1.25, ACCENT)
        _comparison_quadrant(slide, chunk, "BUG", "BUG占比", 6.85, 1.25, SECONDARY)
        _comparison_quadrant(slide, chunk, "氪金", "氪金占比", 0.65, 4.02, AMBER)
        _top_category_panel(slide, chunk, 6.85, 4.02)


def _add_key_insights(
    presentation: Presentation,
    rows: list[dict],
    summary: str,
    results: list[AnalysisResult],
) -> None:
    slide = _content_slide(presentation, "Key Insights", "AI 洞察与业务解释")
    _insight_chart(slide, rows, "负面率", "负面占比", 0.65, 1.22, ACCENT)
    _insight_text(slide, "AI Insight", _findings(summary, results, rows, 3), 7.3, 1.22)
    _insight_chart(slide, rows, "商业化压力", "氪金占比", 0.65, 4.15, AMBER)
    _insight_text(slide, "Market Interpretation", _market_interpretations(rows), 7.3, 4.15)


def _add_representative_reviews(
    presentation: Presentation,
    rows: list[dict],
    results: list[AnalysisResult],
) -> None:
    slide = _content_slide(presentation, "Representative Reviews", "代表性玩家声音")
    reviews_by_market = _reviews_by_market(rows, results)
    markets = list(reviews_by_market.items())[:3]
    if not markets:
        _text(slide, "暂无代表性评论样本", 0.75, 3.1, 11.8, 0.35, 18, MUTED, align=PP_ALIGN.CENTER)
        return
    for index, (market, reviews) in enumerate(markets):
        x = 0.65 + index * 4.12
        _review_market_column(slide, market, reviews, x, 1.25)


def _add_recommendations(presentation: Presentation, source_text: str) -> None:
    slide = _content_slide(presentation, "Actionable Recommendations", "P0 / P1 / P2 优先级行动")
    recommendations = _recommendations(source_text)
    colors = {"P0": ACCENT, "P1": PRIMARY, "P2": SECONDARY}
    for index, level in enumerate(["P0", "P1", "P2"]):
        x = 0.65 + index * 4.12
        _recommendation_column(slide, level, recommendations[level], x, 1.25, colors[level])


def _content_slide(presentation: Presentation, title: str, subtitle: str):
    slide = _blank_slide(presentation)
    _text(slide, title.upper(), 0.65, 0.34, 7.8, 0.42, 29, TEXT, bold=True)
    _text(slide, subtitle, 0.67, 0.82, 6.8, 0.23, 13, MUTED)
    _rect(slide, 0.65, 1.07, 1.18, 0.045, PRIMARY, radius=False, line=PRIMARY)
    _text(slide, "SYSTEM // GAMEPULSE", 9.1, 0.82, 2.1, 0.15, 8, MUTED, bold=True, align=PP_ALIGN.RIGHT)
    _brand_mark(slide, 11.5, 0.38)
    return slide


def _finalize(presentation: Presentation) -> None:
    total = len(presentation.slides)
    for index, slide in enumerate(presentation.slides, start=1):
        _footer(slide, index, total)


def _single_market_row(result: AnalysisResult, context: dict | None = None) -> dict:
    total = max(len(result.classified_reviews), 1)
    sentiments = _sentiment_mix(result)
    return {
        "市场": context.get("market", "当前市场") if context else "当前市场",
        "有效评论数量": len(result.classified_reviews),
        "平均星级": "--",
        "正面占比": _pct_text(sentiments["正面"]),
        "中性占比": _pct_text(sentiments["中性"]),
        "负面占比": _pct_text(sentiments["负面"]),
        "游戏玩法占比": _pct_text(result.category_counts.get("游戏玩法", 0) / total),
        "BUG占比": _pct_text(result.category_counts.get("BUG", 0) / total),
        "氪金占比": _pct_text(result.category_counts.get("氪金", 0) / total),
        "美术占比": _pct_text(result.category_counts.get("美术", 0) / total),
        "性能问题占比": _pct_text(result.category_counts.get("性能优化", 0) / total),
        "整体评价占比": _pct_text(result.category_counts.get("整体评价", 0) / total),
        "其他占比": _pct_text(result.category_counts.get("其他", 0) / total),
    }


def _normalize_context(package_name: str, context: dict | None = None) -> dict:
    context = dict(context or {})
    game_title = str(context.get("game_title") or context.get("game_name") or package_name or "Unknown Game").strip()
    return {
        "game_title": game_title,
        "package_name": str(context.get("package_name") or package_name or "").strip(),
        "market": str(context.get("market") or context.get("scope") or "当前市场").strip(),
        "language": str(context.get("language") or "未记录").strip(),
        "time_range": str(context.get("time_range") or context.get("time_display") or "未记录").strip(),
        "analysis_time": str(context.get("analysis_time") or context.get("analysis_timestamp") or REPORT_TIME).strip(),
    }


def _dashboard_metrics(rows: list[dict]) -> dict[str, str]:
    negative = _avg_percent(rows, "负面占比")
    bug = _avg_percent(rows, "BUG占比")
    monetization = _avg_percent(rows, "氪金占比")
    top_category = _overall_top_category(rows)
    return {
        "negative": f"{negative:.1f}%",
        "bug": f"{bug:.1f}%",
        "monetization": f"{monetization:.1f}%",
        "top_category": top_category,
    }


def _comparison_quadrant(slide, rows: list[dict], title: str, key: str, x: float, y: float, color: RGBColor) -> None:
    _card(slide, x, y, 5.8, 2.35)
    _text(slide, title, x + 0.25, y + 0.18, 2.2, 0.26, 17, TEXT, bold=True)
    values = [(str(row.get("市场", "")), _parse_percent(row.get(key, 0))) for row in rows]
    _mini_bars(slide, values, x + 0.28, y + 0.68, 5.15, 1.35, color)


def _top_category_panel(slide, rows: list[dict], x: float, y: float) -> None:
    _card(slide, x, y, 5.8, 2.35)
    _text(slide, "Top Category", x + 0.25, y + 0.18, 2.2, 0.26, 17, TEXT, bold=True)
    for index, row in enumerate(rows[:4]):
        category, value = _row_top_category(row)
        yy = y + 0.68 + index * 0.34
        _text(slide, str(row.get("市场", "")), x + 0.28, yy, 1.25, 0.2, 10, MUTED, bold=True)
        _text(slide, category, x + 1.62, yy, 1.45, 0.2, 11, TEXT, bold=True)
        _text(slide, f"{value:.1f}%", x + 4.65, yy, 0.72, 0.2, 11, AMBER, bold=True, align=PP_ALIGN.RIGHT)


def _insight_chart(slide, rows: list[dict], title: str, key: str, x: float, y: float, color: RGBColor) -> None:
    _card(slide, x, y, 6.2, 2.35)
    _text(slide, title, x + 0.25, y + 0.2, 2.2, 0.25, 17, TEXT, bold=True)
    values = [(str(row.get("市场", "")), _parse_percent(row.get(key, 0))) for row in rows[:5]]
    _mini_bars(slide, values, x + 0.3, y + 0.72, 5.45, 1.22, color)


def _insight_text(slide, title: str, items: list[str], x: float, y: float) -> None:
    _card(slide, x, y, 5.35, 2.35)
    _text(slide, title, x + 0.25, y + 0.2, 2.3, 0.25, 17, TEXT, bold=True)
    for index, item in enumerate((items or ["暂无洞察"])[:3], start=1):
        _number_dot(slide, index, x + 0.28, y + 0.66 + (index - 1) * 0.47)
        _text(slide, item, x + 0.68, y + 0.61 + (index - 1) * 0.47, 4.25, 0.32, 12, TEXT)


def _review_market_column(slide, market: str, reviews: list, x: float, y: float) -> None:
    _card(slide, x, y, 3.75, 5.55)
    _text(slide, market, x + 0.25, y + 0.22, 2.8, 0.28, 18, TEXT, bold=True)
    positive = next((review for review in reviews if review.sentiment == "正面"), reviews[0])
    negative = next((review for review in reviews if review.sentiment == "负面"), reviews[-1])
    _review_block(slide, "Positive", positive, x + 0.25, y + 0.75, GREEN)
    _review_block(slide, "Negative", negative, x + 0.25, y + 3.08, ACCENT)


def _review_block(slide, label: str, review, x: float, y: float, color: RGBColor) -> None:
    label_color = _rgb("#FFFFFF") if color == ACCENT else color
    _text(slide, label, x, y, 1.2, 0.2, 11, label_color, bold=True)
    _text(slide, _truncate(review.content, 72), x, y + 0.35, 3.15, 0.5, 10, TEXT)
    _text(slide, "↓", x + 1.45, y + 0.93, 0.25, 0.16, 10, MUTED, align=PP_ALIGN.CENTER)
    _text(slide, _truncate(review.reason or "暂无中文摘要", 50), x, y + 1.2, 3.15, 0.36, 10, MUTED)
    _pill(slide, f"{review.category} · {review.sentiment}", x, y + 1.72, color)


def _recommendation_column(slide, level: str, items: list[dict], x: float, y: float, color: RGBColor) -> None:
    _card(slide, x, y, 3.75, 5.55)
    label = {"P0": "立即优化", "P1": "建议优化", "P2": "长期优化"}[level]
    _rect(slide, x + 0.25, y + 0.22, 0.58, 0.36, color, line=color)
    label_color = _contrast_text_rgb(color)
    _text(slide, level, x + 0.25, y + 0.29, 0.58, 0.12, 12, label_color, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, label, x + 0.95, y + 0.25, 2.1, 0.25, 17, TEXT, bold=True)
    for index, item in enumerate(items[:2]):
        yy = y + 0.95 + index * 2.05
        _text(slide, item["title"], x + 0.28, yy, 3.15, 0.24, 13, TEXT, bold=True)
        _text(slide, "原因：" + item["reason"], x + 0.28, yy + 0.45, 3.1, 0.35, 10, MUTED)
        _text(slide, "收益：" + item["impact"], x + 0.28, yy + 1.08, 3.1, 0.35, 10, MUTED)


def _cover_info(slide, label: str, value: str, x: float, y: float) -> None:
    _text(slide, label.upper(), x, y, 1.65, 0.21, 10, TEXT, bold=True)
    _rect(slide, x, y + 0.24, 0.42, 0.035, PRIMARY, radius=False, line=PRIMARY)
    _text(slide, _truncate(value, 55), x + 1.82, y - 0.02, 5.2, 0.25, 13, TEXT, bold=True)


def _mini_dashboard(slide, x: float, y: float) -> None:
    for index, label in enumerate(["Sentiment", "Category", "Risk", "Action"]):
        xx = x + (index % 2) * 1.65
        yy = y + (index // 2) * 1.15
        _rect(slide, xx, yy, 1.28, 0.78, NAVY, transparency=0, line=BORDER)
        _corner_brackets(slide, xx, yy, 1.28, 0.78, PRIMARY if index == 0 else SECONDARY)
        _text(slide, label.upper(), xx + 0.15, yy + 0.2, 0.98, 0.18, 9, TEXT, bold=True, align=PP_ALIGN.CENTER)


def _dashboard_tile(slide, title: str, value: str, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    _card(slide, x, y, w, h)
    _text(slide, title.upper(), x + 0.28, y + 0.24, 2.2, 0.26, 13, MUTED, bold=True)
    value_color = AMBER if color == PRIMARY else color
    _text(slide, value, x + 0.28, y + 0.72, 2.8, 0.52, 31, value_color, bold=True)
    _rect(slide, x + w - 1.35, y + 0.45, 0.86, 0.86, color, transparency=55, line=color)
    _text(slide, "●", x + w - 1.13, y + 0.62, 0.42, 0.22, 16, _contrast_text_rgb(color), bold=True, align=PP_ALIGN.CENTER)


def _card(slide, x: float, y: float, w: float, h: float) -> None:
    _rect(slide, x, y, w, h, CARD, radius=False, line=BORDER)
    _corner_brackets(slide, x, y, w, h, PRIMARY)


def add_card(slide, x: float, y: float, w: float, h: float):
    _card(slide, x, y, w, h)


def add_badge(slide, text: str, x: float, y: float, fill: RGBColor = PRIMARY) -> None:
    _rect(slide, x, y, 0.72, 0.28, fill, line=fill)
    _text(slide, text, x + 0.04, y + 0.065, 0.64, 0.08, 8, _contrast_text_rgb(fill), bold=True, align=PP_ALIGN.CENTER)


def _rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    radius: bool = True,
    transparency: int = 0,
    line: RGBColor | None = None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line or fill
    return shape


def _hud_grid(slide) -> None:
    for x in [0.7, 2.1, 3.5, 4.9, 6.3, 7.7, 9.1, 10.5, 11.9]:
        _rect(slide, x, 0.35, 0.006, 6.55, GRID, radius=False, transparency=25, line=GRID)
    for y in [1.1, 2.2, 3.3, 4.4, 5.5, 6.6]:
        _rect(slide, 0.45, y, 12.35, 0.006, GRID, radius=False, transparency=25, line=GRID)


def _corner_brackets(slide, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    length = 0.22
    thickness = 0.018
    for dx, dy, hx, hy in [
        (0, 0, length, thickness), (0, 0, thickness, length),
        (w - length, 0, length, thickness), (w - thickness, 0, thickness, length),
        (0, h - thickness, length, thickness), (0, h - length, thickness, length),
        (w - length, h - thickness, length, thickness), (w - thickness, h - length, thickness, length),
    ]:
        _rect(slide, x + dx, y + dy, hx, hy, color, radius=False, line=color)


def _text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for attr in ["margin_left", "margin_right", "margin_top", "margin_bottom"]:
        setattr(frame, attr, 0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = _clean_markdown(text)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    background_color: RGBColor = CARD,
    text_color: RGBColor | None = None,
    bold: bool = False,
) -> None:
    _rect(slide, x, y, w, h, background_color, line=BORDER)
    _text(slide, text, x + 0.08, y + 0.08, max(w - 0.16, 0.1), max(h - 0.16, 0.1), size, text_color or _contrast_text_rgb(background_color), bold=bold)


def _mini_bars(slide, values: list[tuple[str, float]], x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    if not values:
        _text(slide, "暂无数据", x, y + 0.5, w, 0.22, 12, MUTED, align=PP_ALIGN.CENTER)
        return
    max_value = max(value for _, value in values) or 1
    gap = h / max(len(values), 1)
    label_w = 1.4
    value_w = 0.7
    bar_w = w - label_w - value_w - 0.25
    for index, (label, value) in enumerate(values[:6]):
        yy = y + index * gap
        _text(slide, _truncate(label, 12), x, yy + 0.04, label_w, 0.16, 9, MUTED, bold=True)
        _rect(slide, x + label_w, yy + 0.08, bar_w, 0.12, BORDER, radius=False, line=BORDER)
        _rect(slide, x + label_w, yy + 0.08, bar_w * value / max_value, 0.12, color, radius=False, line=color)
        _text(slide, f"{value:.1f}%", x + label_w + bar_w + 0.08, yy + 0.03, value_w, 0.16, 9, TEXT, bold=True, align=PP_ALIGN.RIGHT)


def _bullet_cards(slide, items: list[dict], x: float, y: float, w: float, h: float, columns: int = 1) -> None:
    if not items:
        items = [{"title": "暂无足够信息", "detail": "当前样本不足以形成稳定结论。"}]
    col_w = w / columns - 0.12
    row_h = min(1.02, h / max((len(items) + columns - 1) // columns, 1) - 0.08)
    for index, item in enumerate(items[:6]):
        col = index % columns
        row = index // columns
        xx = x + col * (col_w + 0.24)
        yy = y + row * (row_h + 0.16)
        _card(slide, xx, yy, col_w, row_h)
        _rect(slide, xx + 0.16, yy + 0.17, 0.08, 0.5, PRIMARY, radius=False, line=PRIMARY)
        _text(slide, item.get("title", "要点"), xx + 0.34, yy + 0.15, col_w - 0.52, 0.18, 11, TEXT, bold=True)
        _text(slide, item.get("detail", ""), xx + 0.34, yy + 0.43, col_w - 0.52, max(row_h - 0.55, 0.25), 9, MUTED)


def _number_dot(slide, number: int, x: float, y: float) -> None:
    _rect(slide, x, y, 0.24, 0.24, PRIMARY, line=PRIMARY)
    _text(slide, str(number), x, y + 0.045, 0.24, 0.08, 7, TEXT, bold=True, align=PP_ALIGN.CENTER)


def _pill(slide, text: str, x: float, y: float, color: RGBColor) -> None:
    _rect(slide, x, y, 1.72, 0.26, color, transparency=45, line=color)
    _text(slide, text, x + 0.08, y + 0.06, 1.55, 0.08, 8, _contrast_text_rgb(color), bold=True, align=PP_ALIGN.CENTER)


def _brand_mark(slide, x: float, y: float) -> None:
    _triangle_logo(slide, x, y)
    _text(slide, "Market Intelligence", x + 0.35, y + 0.02, 1.35, 0.14, 8, MUTED, bold=True)


def _triangle_logo(slide, x: float, y: float) -> None:
    _rect(slide, x, y + 0.02, 0.18, 0.28, PRIMARY, line=PRIMARY)
    _rect(slide, x + 0.12, y + 0.08, 0.18, 0.22, GREEN, line=GREEN)
    _rect(slide, x + 0.21, y + 0.14, 0.14, 0.14, SECONDARY, line=SECONDARY)


def _footer(slide, page: int, total: int) -> None:
    _text(slide, "GAMEPULSE // MARKET INTELLIGENCE", 0.65, 7.04, 3.6, 0.14, 8, MUTED)
    _text(slide, REPORT_TIME, 5.3, 7.04, 2.7, 0.14, 8, MUTED, align=PP_ALIGN.CENTER)
    _text(slide, f"{page}/{total}", 11.95, 7.04, 0.75, 0.14, 8, MUTED, align=PP_ALIGN.RIGHT)


def _total_comments(rows: list[dict], results: list[AnalysisResult]) -> int:
    row_total = sum(_int_value(row.get("有效评论数量", 0)) for row in rows)
    if row_total:
        return row_total
    return sum(len(result.classified_reviews) for result in results)


def _average_score(rows: list[dict]) -> str:
    values = [_float_value(row.get("平均星级")) for row in rows]
    values = [value for value in values if value is not None]
    return f"{mean(values):.2f}" if values else "--"


def _sentiment_mix(result: AnalysisResult) -> dict[str, float]:
    total = max(len(result.classified_reviews), 1)
    return {
        sentiment: sum(1 for review in result.classified_reviews if review.sentiment == sentiment) / total
        for sentiment in ["正面", "中性", "负面"]
    }


def _reviews_by_market(rows: list[dict], results: list[AnalysisResult]) -> dict[str, list]:
    markets = [str(row.get("市场", f"市场{index + 1}")) for index, row in enumerate(rows)]
    output = {}
    for index, result in enumerate(results[: len(markets)]):
        output[markets[index]] = result.classified_reviews[:4]
    if not output and results:
        output[markets[0] if markets else "当前市场"] = results[0].classified_reviews[:4]
    return {market: reviews for market, reviews in output.items() if reviews}


def _findings(summary: str, results: list[AnalysisResult], rows: list[dict], limit: int) -> list[str]:
    lines = _extract_lines(summary, 12)
    if lines:
        return lines[:limit]
    findings = []
    if rows:
        highest_negative = max(rows, key=lambda row: _parse_percent(row.get("负面占比", 0)))
        findings.append(f"{highest_negative.get('市场', '重点市场')}负面率相对突出，需要优先复盘差评样本")
        findings.append(f"主要高频类别集中在{_overall_top_category(rows)}，建议作为近期运营观察重点")
    for result in results[:1]:
        findings.extend(result.most_satisfied[:1])
        findings.extend(result.most_unsatisfied[:1])
    return (findings + ["样本反馈显示口碑表现存在结构性差异"])[:limit]


def _market_interpretations(rows: list[dict]) -> list[str]:
    if not rows:
        return ["暂无市场样本"]
    high_bug = max(rows, key=lambda row: _parse_percent(row.get("BUG占比", 0)))
    high_pay = max(rows, key=lambda row: _parse_percent(row.get("氪金占比", 0)))
    high_positive = max(rows, key=lambda row: _parse_percent(row.get("正面占比", 0)))
    return [
        f"{high_positive.get('市场', '优势市场')}正面反馈更集中，可提炼为素材卖点",
        f"{high_bug.get('市场', '重点市场')}BUG反馈相对突出，需关注稳定性风险",
        f"{high_pay.get('市场', '重点市场')}商业化争议更高，建议复核付费体验",
    ]


def _management_summary(summary: str, rows: list[dict]) -> str:
    lines = _extract_lines(summary, 1)
    if lines:
        return _truncate(lines[0], 86)
    top = _overall_top_category(rows)
    negative = _avg_percent(rows, "负面占比")
    return f"本次样本显示主要反馈集中在{top}，平均负面率约{negative:.1f}%，建议优先处理跨市场共性痛点。"


def _recommendations(source_text: str) -> dict[str, list[dict]]:
    lines = _extract_lines(source_text, 8)
    defaults = [
        "优先处理跨市场高频负面问题",
        "复盘商业化争议与付费体验",
        "建立版本稳定性专项监控",
        "提炼正向反馈用于买量素材",
        "按区域补充样本验证本地化风险",
        "持续追踪活动运营口碑变化",
    ]
    items = (lines + defaults)[:6]
    return {
        "P0": [
            {"title": _truncate(items[0], 18), "reason": "多个样本反馈显示该问题影响核心体验", "impact": "降低差评风险，改善短期口碑"},
            {"title": _truncate(items[1], 18), "reason": "商业化与体验冲突容易放大负面情绪", "impact": "提升付费接受度与留存质量"},
        ],
        "P1": [
            {"title": _truncate(items[2], 18), "reason": "稳定性问题会持续影响版本评分", "impact": "改善用户信任与商店评分"},
            {"title": _truncate(items[3], 18), "reason": "正向反馈可转化为市场沟通资产", "impact": "提升素材相关性与转化效率"},
        ],
        "P2": [
            {"title": _truncate(items[4], 18), "reason": "区域差异需要更长期的本地化验证", "impact": "降低发行和运营决策偏差"},
            {"title": _truncate(items[5], 18), "reason": "活动口碑变化能反映运营节奏质量", "impact": "优化长期内容规划"},
        ],
    }


def _overall_top_category(rows: list[dict]) -> str:
    keys = ["BUG占比", "氪金占比", "性能问题占比", "游戏玩法占比", "美术占比", "整体评价占比", "其他占比"]
    totals = {key: sum(_parse_percent(row.get(key, 0)) for row in rows) for key in keys}
    return max(totals, key=totals.get).replace("占比", "")


def _row_top_category(row: dict) -> tuple[str, float]:
    keys = ["BUG占比", "氪金占比", "性能问题占比", "游戏玩法占比", "美术占比", "整体评价占比", "其他占比"]
    key = max(keys, key=lambda item: _parse_percent(row.get(item, 0)))
    return key.replace("占比", ""), _parse_percent(row.get(key, 0))


def _avg_percent(rows: list[dict], key: str) -> float:
    if not rows:
        return 0.0
    return mean(_parse_percent(row.get(key, 0)) for row in rows)


def _extract_lines(text: str, limit: int) -> list[str]:
    cleaned = _clean_markdown(text)
    parts = re.split(r"[\n。；;]+", cleaned)
    return [_truncate(part.strip(), 56) for part in parts if part.strip()][:limit]


def _clean_markdown(text: str) -> str:
    value = str(text or "")
    value = re.sub(r"```.*?```", "", value, flags=re.S)
    value = re.sub(r"^\s*#{1,6}\s*", "", value, flags=re.M)
    value = re.sub(r"\*\*(.*?)\*\*", r"\1", value)
    value = re.sub(r"\*(.*?)\*", r"\1", value)
    value = re.sub(r"`([^`]*)`", r"\1", value)
    value = re.sub(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$", "", value, flags=re.M)
    value = re.sub(r"^\s*[-*]\s+", "", value, flags=re.M)
    value = re.sub(r"^\s*\d+[.)]\s+", "", value, flags=re.M)
    value = value.replace("|", " ")
    value = value.replace("---", " ")
    return re.sub(r"\s+", " ", value).strip()


def _truncate(text: str, limit: int) -> str:
    cleaned = _clean_markdown(text)
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1] + "…"


def _parse_percent(value) -> float:
    try:
        return float(str(value).replace("%", "").strip())
    except ValueError:
        return 0.0


def _float_value(value) -> float | None:
    try:
        return float(str(value).replace("%", "").strip())
    except (TypeError, ValueError):
        return None


def _int_value(value) -> int:
    try:
        return int(float(str(value).replace("%", "").strip()))
    except (TypeError, ValueError):
        return 0


def _pct_text(value: float) -> str:
    return f"{value * 100:.1f}%"
